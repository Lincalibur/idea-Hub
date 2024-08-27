import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def convert_svg_unit(value, total_dimension=None):
    """Convert SVG units to inches for PowerPoint. Handles both absolute and percentage values."""
    if value.endswith('%'):
        if total_dimension is not None:
            return Inches((float(value.strip('%')) / 100.0) * total_dimension)
        else:
            raise ValueError("Percentage value provided without a reference dimension.")
    else:
        return Inches(float(value) / 96)  # Convert pixels to inches (assuming 96 dpi)

def parse_points(point_string, total_width, total_height):
    """Parse a list of points from an SVG 'points' attribute."""
    points = []
    for point in point_string.strip().split():
        x, y = map(float, point.split(','))
        points.append((convert_svg_unit(str(x), total_width), convert_svg_unit(str(y), total_height)))
    return points

def parse_path_data(path_data, total_width, total_height):
    """Parse an SVG path's data attribute for M (move to) and L (line to) commands."""
    points = []
    commands = path_data.strip().split()
    
    for command in commands:
        if command.startswith(('M', 'L')):
            coord_pairs = command[1:].strip().split(' ')
            for pair in coord_pairs:
                if ',' in pair:
                    x, y = map(float, pair.split(','))
                    points.append((convert_svg_unit(str(x), total_width), convert_svg_unit(str(y), total_height)))
        elif ',' in command:
            x, y = map(float, command.split(','))
            points.append((convert_svg_unit(str(x), total_width), convert_svg_unit(str(y), total_height)))
    return points

# Path to the SVG file
svg_file_path = 'Converters\Draw To Power Point\TestDiagram.drawio.svg'

# Define the total dimensions of the SVG (in pixels or any other unit you prefer)
total_width = 800  # Example total width of the SVG canvas
total_height = 600  # Example total height of the SVG canvas

# Create a new PowerPoint presentation
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Choosing a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Parse the SVG file
tree = ET.parse(svg_file_path)
root = tree.getroot()

# Namespace dictionary to handle SVG's default namespace
ns = {'svg': 'http://www.w3.org/2000/svg'}

# Loop through SVG elements and add them to the PowerPoint slide
for elem in root.findall('.//svg:rect', ns):
    x = convert_svg_unit(elem.attrib['x'], total_width)
    y = convert_svg_unit(elem.attrib['y'], total_height)
    width = convert_svg_unit(elem.attrib['width'])
    height = convert_svg_unit(elem.attrib['height'])
    
    # Add a rectangle shape (representing process or similar) to the slide
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    
    # Look for text elements inside the rectangle
    for text_elem in root.findall('.//svg:text', ns):
        text_x = convert_svg_unit(text_elem.attrib['x'], total_width)
        text_y = convert_svg_unit(text_elem.attrib['y'], total_height)
        text_content = text_elem.text.strip()

        # Check if the text is within the bounds of the rectangle
        if x <= text_x <= x + width and y <= text_y <= y + height:
            # Center the text inside the shape
            text_box = shape.text_frame
            text_box.text = text_content
            text_box.paragraphs[0].font.size = Pt(12)  # Adjust font size as necessary

for elem in root.findall('.//svg:circle', ns):
    cx = convert_svg_unit(elem.attrib['cx'], total_width)
    cy = convert_svg_unit(elem.attrib['cy'], total_height)
    r = convert_svg_unit(elem.attrib['r'])
    
    # Add an oval shape (representing terminator or similar) to the slide
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r)

    # Look for text elements inside the circle
    for text_elem in root.findall('.//svg:text', ns):
        text_x = convert_svg_unit(text_elem.attrib['x'], total_width)
        text_y = convert_svg_unit(text_elem.attrib['y'], total_height)
        text_content = text_elem.text.strip()

        # Check if the text is within the bounds of the circle
        if (cx - r) <= text_x <= (cx + r) and (cy - r) <= (cy + r):
            # Center the text inside the shape
            text_box = shape.text_frame
            text_box.text = text_content
            text_box.paragraphs[0].font.size = Pt(12)  # Adjust font size as necessary

# Handle connectors (lines)
for elem in root.findall('.//svg:line', ns):
    x1 = convert_svg_unit(elem.attrib['x1'], total_width)
    y1 = convert_svg_unit(elem.attrib['y1'], total_height)
    x2 = convert_svg_unit(elem.attrib['x2'], total_width)
    y2 = convert_svg_unit(elem.attrib['y2'], total_height)
    
    # Add a line shape to the slide
    connector = slide.shapes.add_connector(MSO_SHAPE.LINE, x1, y1, x2, y2)
    connector.line.color.rgb = RGBColor(0, 0, 0)  # Black color

# Handle connectors (polylines)
for elem in root.findall('.//svg:polyline', ns):
    points = parse_points(elem.attrib['points'], total_width, total_height)
    
    # Draw lines between consecutive points
    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        connector = slide.shapes.add_connector(MSO_SHAPE.LINE, x1, y1, x2, y2)
        connector.line.color.rgb = RGBColor(0, 0, 0)  # Black color

# Handle connectors (paths with 'M' and 'L' commands for simplicity)
for elem in root.findall('.//svg:path', ns):
    path_data = elem.attrib['d']
    points = parse_path_data(path_data, total_width, total_height)
    
    # Draw lines between consecutive points
    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        connector = slide.shapes.add_connector(MSO_SHAPE.LINE, x1, y1, x2, y2)
        connector.line.color.rgb = RGBColor(0, 0, 0)  # Black color

# Save the PowerPoint presentation
output_pptx_path = 'Converters\Draw To Power Point\Diagram.pptx'
presentation.save(output_pptx_path)

print(f"PowerPoint presentation saved as {output_pptx_path}")
